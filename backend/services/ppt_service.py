"""PPT / PDF text extraction + LLM summarisation."""
from __future__ import annotations

import base64
import logging
import os
import tempfile

log = logging.getLogger(__name__)


def extract_and_summarize(payload: dict) -> dict:
    file_b64 = payload.get("file_b64", "")
    file_type = payload.get("file_type", "pdf").lower().lstrip(".")
    title = payload.get("title") or "Uploaded File"

    if not file_b64:
        return {"title": title, "text": "No file data provided.", "summary": ""}

    file_data = base64.b64decode(file_b64)
    suffix = f".{file_type}"

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(file_data)
            tmp_path = f.name

        if file_type in ("pptx", "ppt"):
            text = _extract_pptx(tmp_path)
        elif file_type == "pdf":
            text = _extract_pdf(tmp_path)
        else:
            text = f"Unsupported file type: {file_type}. Supported: pdf, pptx."
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    summary = _summarize(text, title)
    return {"title": f"{title} Analysis", "text": text, "summary": summary}


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        lines: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines) if lines else "No readable text found in PPTX."
    except ImportError:
        return "python-pptx not installed. Run: pip install python-pptx"
    except Exception as exc:
        return f"PPTX extraction failed: {exc}"


def _extract_pdf(path: str) -> str:
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        pages: list[str] = []
        for i, page in enumerate(doc, 1):
            page_text = page.get_text().strip()
            if page_text:
                pages.append(f"[Page {i}]\n{page_text}")
        doc.close()
        return "\n\n".join(pages) if pages else "No readable text found in PDF."
    except ImportError:
        return "PyMuPDF not installed. Run: pip install PyMuPDF"
    except Exception as exc:
        return f"PDF extraction failed: {exc}"


def _summarize(text: str, title: str) -> str:
    if not text or text.startswith(("No readable", "Unsupported", "python-pptx", "PyMuPDF")):
        return text

    try:
        from config import CHAT_MODEL
        from services.ai_service import _client, _track_usage

        prompt = (
            f"You are an academic study assistant. Below is content extracted from a file titled '{title}'.\n\n"
            f"CONTENT:\n{text[:6000]}\n\n"
            "Please provide:\n"
            "1. Key Topics Outline\n"
            "2. Important Terms & Definitions\n"
            "3. Exam Hotspots / Key Takeaways\n\n"
            "Format with clear headers."
        )
        resp = _client().chat.completions.create(
            model=CHAT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500,
        )
        summary = resp.choices[0].message.content
        _track_usage(CHAT_MODEL, resp.usage.prompt_tokens, resp.usage.completion_tokens, "ppt_analyze")
        return summary
    except Exception as exc:
        log.error("PPT summarise failed: %s", exc)
        return f"AI summary unavailable. Error: {exc}"
